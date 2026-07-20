from math import ceil
import os
from pathlib import Path
from typing import Optional
from pptx import Presentation
import fileindexer_extract as _native

from config import Settings

class FileProcessor:
    """Class to handle file processing and text extraction."""
    settings = Settings()

    @staticmethod
    def extract_text_from_pdf(file_path: str) -> Optional[str]:
        """Extract text from a PDF file (Rust implementation)."""
        return _native.extract_text_from_pdf(str(file_path))

    @staticmethod
    def extract_text_from_docx(file_path: str) -> Optional[str]:
        """Extract text from a DOCX file (Rust implementation)."""
        return _native.extract_text_from_docx(str(file_path))
    
    @staticmethod
    def extract_text_from_pptx(file_path: str) -> Optional[str]:
        """Extract text from a PPTX file."""
        try:
            prs = Presentation(file_path)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return "\n".join(text)
        except Exception as e:
            print(f"Error extracting text from PPTX: {e}")
            return None
    
    @staticmethod
    def extract_text(file_path: str) -> Optional[str]:
        """Extract text from a TXT or MD file (Rust implementation)."""
        return _native.extract_text(str(file_path))

    @staticmethod
    def process_file(file_path: str) -> Optional[str]:
        """Process a file and extract its text based on the file type."""
        extension = Path(file_path).suffix.lower()
        if extension == '.pptx':
            return FileProcessor.extract_text_from_pptx(file_path)
        return _native.process_file(str(file_path))

    @staticmethod
    def process_files_parallel(file_paths: list) -> list:
        """Extract text from every file in parallel (Rust, GIL released).

        Assumes no .pptx paths are present (not in VALID_FILE_EXTENSIONS,
        and .pptx isn't supported by the Rust extractor).
        """
        return _native.process_files_parallel([str(p) for p in file_paths])
    
    @staticmethod
    def chunk_text(text: str, chunk_size: int = settings.CHUNK_SIZE, overlap: int = settings.CHUNK_OVERLAP) -> list:
        """Chunk text into smaller pieces."""
        chunks = []
        start = 0
        text_length = len(text)

        if (text_length <= chunk_size):
            return [text]

        num_chunks = ceil((text_length - overlap) / (chunk_size - overlap))
        chunk_size = ceil(text_length / num_chunks) + overlap

        
        while start < text_length:
            end = min(start + chunk_size, text_length)
            chunks.append(text[start:end])
            start += chunk_size - overlap
            
        return chunks